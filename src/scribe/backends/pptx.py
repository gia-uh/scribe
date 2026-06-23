from __future__ import annotations

import io

from pptx import Presentation

from ..result import ExtractResult


def convert(data: bytes) -> ExtractResult:
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ValueError(f"could not open PPTX: {exc}") from exc

    parts: list[str] = []
    title: str | None = None
    slides = list(prs.slides)
    for i, slide in enumerate(slides, 1):
        stitle = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            stitle = slide.shapes.title.text.strip()
        if title is None and stitle:
            title = stitle
        parts.append(f"## Slide {i}" + (f": {stitle}" if stitle else ""))
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip() or para.text.strip()
                    if text:
                        parts.append("- " + text)
            if shape.has_table:
                tbl = shape.table
                rows = [[c.text.strip() for c in row.cells] for row in tbl.rows]
                if rows:
                    head, *body = rows
                    parts.append("| " + " | ".join(head) + " |")
                    parts.append("| " + " | ".join(["---"] * len(head)) + " |")
                    for r in body:
                        parts.append("| " + " | ".join(r) + " |")
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append("> Notes: " + note.replace("\n", " "))

    return ExtractResult(
        markdown="\n\n".join(parts), title=title, meta={"backend": "pptx", "slides": len(slides)}
    )
