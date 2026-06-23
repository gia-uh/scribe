import io

from pptx import Presentation

from scribe.backends.pptx import convert


def _make():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Intro"
    slide.placeholders[1].text = "point one\npoint two"
    slide.notes_slide.notes_text_frame.text = "speaker note here"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_slides_bullets_notes():
    md = convert(_make()).markdown
    assert "Intro" in md
    assert "- point one" in md and "- point two" in md
    assert "speaker note here" in md


def test_pptx_meta_slides():
    assert convert(_make()).meta["slides"] == 1


def test_pptx_corrupt_raises():
    import pytest

    with pytest.raises(ValueError):
        convert(b"not a pptx")
